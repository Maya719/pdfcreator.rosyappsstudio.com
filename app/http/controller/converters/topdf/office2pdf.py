import asyncio
import csv as csv_module
from io import BytesIO, StringIO
from fastapi import HTTPException
from app.helpers.pdf_utils import (
    _make_output_dir, _html_to_pdf_bytes, _build_html_table, _wrap_html, 
    _rl_text_to_pdf, _extract_odf_text, APP_URL
)

# ReportLab imports for PPTX
from reportlab.platypus import Paragraph, Spacer, SimpleDocTemplate
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import mm

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    from odf.opendocument import load as odf_load
    from odf import table as odf_table
except ImportError:
    odf_load = None

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None

class OfficeToPdfConverter:
    @staticmethod
    async def convert_csv(content: bytes, filename: str):
        output_dir = _make_output_dir()
        output_path = output_dir / filename
        try:
            text = content.decode("utf-8", errors="replace")
            reader = csv_module.reader(StringIO(text))
            rows = list(reader)
            if not rows: raise Exception("CSV is empty")
            html = _wrap_html(_build_html_table(rows[0], rows[1:], "CSV Data"))
            await _html_to_pdf_bytes(html, output_path)
            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"CSV to PDF failed: {str(e)}")

    @staticmethod
    async def convert_xlsx(content: bytes, filename: str):
        if not openpyxl: raise HTTPException(status_code=500, detail="openpyxl not installed")
        output_dir = _make_output_dir()
        output_path = output_dir / filename
        try:
            wb = openpyxl.load_workbook(BytesIO(content), data_only=True)
            body = ""
            for name in wb.sheetnames:
                ws = wb[name]
                rows = list(ws.iter_rows(values_only=True))
                if not rows: continue
                body += _build_html_table([str(v) for v in rows[0]], [[str(v) for v in r] for r in rows[1:]], name)
            await _html_to_pdf_bytes(_wrap_html(body), output_path)
            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"XLSX to PDF failed: {str(e)}")

    @staticmethod
    async def convert_pptx(content: bytes, filename: str):
        if not PptxPresentation: raise HTTPException(status_code=500, detail="python-pptx not installed")
        output_dir = _make_output_dir()
        output_path = output_dir / filename
        try:
            def convert():
                prs = PptxPresentation(BytesIO(content))
                doc = SimpleDocTemplate(str(output_path), pagesize=A4)
                styles = getSampleStyleSheet()
                story = []
                for i, slide in enumerate(prs.slides, 1):
                    story.append(Paragraph(f"Slide {i}", styles["Heading1"]))
                    for shape in slide.shapes:
                        if not shape.has_text_frame: continue
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                story.append(Paragraph(para.text, styles["Normal"]))
                    story.append(Spacer(1, 8 * mm))
                doc.build(story)
            await asyncio.to_thread(convert)
            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PPTX to PDF failed: {str(e)}")
