import uuid
import time
import os
import csv as csv_module
import zipfile
import asyncio
from io import BytesIO, StringIO
from typing import List, Optional, Union
from fastapi import HTTPException, UploadFile

# pyrefly: ignore [missing-import]
from PIL import Image
# pyrefly: ignore [missing-import]
from xhtml2pdf import pisa
from pdf2docx import Converter
import requests

# pyrefly: ignore [missing-import]
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import mm
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Preformatted
)
from reportlab.lib.enums import TA_LEFT

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import mammoth
except ImportError:
    mammoth = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches
except ImportError:
    PptxPresentation = None

try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None

try:
    from odf.opendocument import load as odf_load
    from odf import text as odf_text, table as odf_table
    from odf.element import Element
except ImportError:
    odf_load = None

from config.paths import PRIVATE_DISK, PUBLIC_DISK
from config.app import APP_URL


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _make_output_dir():
    out = PUBLIC_DISK / "converted"
    out.mkdir(parents=True, exist_ok=True)
    return out


def _html_to_pdf_bytes(html: str, output_path):
    """Convert an HTML string to a PDF file at output_path using xhtml2pdf."""
    with open(output_path, "wb") as f:
        status = pisa.CreatePDF(html, dest=f)
    return status


def _build_html_table(headers: list, rows: list, sheet_name: str = None) -> str:
    """Build a minimal HTML table string from headers + rows (lists of values)."""
    head = f"<h2>{sheet_name}</h2>" if sheet_name else ""
    th = "".join(f"<th style='background:#4a90d9;color:#fff;padding:4px 8px;border:1px solid #ccc'>{h}</th>" for h in headers)
    body_rows = ""
    for i, row in enumerate(rows):
        bg = "#f5f5f5" if i % 2 == 0 else "#ffffff"
        tds = "".join(f"<td style='padding:3px 8px;border:1px solid #ccc'>{v}</td>" for v in row)
        body_rows += f"<tr style='background:{bg}'>{tds}</tr>"
    return f"""{head}
<table style='border-collapse:collapse;width:100%;font-size:11px;font-family:Arial,sans-serif'>
  <thead><tr>{th}</tr></thead>
  <tbody>{body_rows}</tbody>
</table>"""


def _wrap_html(body: str) -> str:
    return f"""<!DOCTYPE html>
<html><head>
<meta charset='utf-8'>
<style>
  body {{ font-family: Arial, sans-serif; font-size: 12px; margin: 20px; }}
  h2 {{ color: #333; font-size: 14px; margin-top: 20px; }}
  table {{ page-break-inside: auto; }}
  tr {{ page-break-inside: avoid; }}
</style>
</head><body>{body}</body></html>"""


def _rl_text_to_pdf(text: str, output_path, title: str = "Document"):
    """Render plain text to a PDF using reportlab."""
    doc = SimpleDocTemplate(
        str(output_path),
        pagesize=A4,
        rightMargin=20 * mm,
        leftMargin=20 * mm,
        topMargin=20 * mm,
        bottomMargin=20 * mm,
    )
    styles = getSampleStyleSheet()
    pre_style = ParagraphStyle(
        "pre",
        fontName="Courier",
        fontSize=9,
        leading=13,
        wordWrap="LTR",
        alignment=TA_LEFT,
    )
    story = [Preformatted(text or "(empty document)", pre_style)]
    doc.build(story)


# ---------------------------------------------------------------------------
# ODF text extraction helper
# ---------------------------------------------------------------------------

def _extract_odf_text(odf_doc) -> str:
    """Extract all text from an ODF document (odt/ods) as plain text."""
    lines = []

    def recurse(element):
        if hasattr(element, "childNodes"):
            for child in element.childNodes:
                if hasattr(child, "data"):
                    lines.append(child.data)
                elif hasattr(child, "tagName"):
                    tag = child.tagName.lower() if child.tagName else ""
                    if "table-row" in tag or "p" in tag:
                        lines.append("\n")
                    recurse(child)

    recurse(odf_doc.body)
    return "".join(lines)


class UniversalConverterController:

    # -----------------------------------------------------------------------
    # Password-protection validator
    # -----------------------------------------------------------------------

    @staticmethod
    def _validate_file_not_protected(file_bytes: bytes, ext: str):
        """
        Checks if a file is password protected using its raw bytes.
        """
        if ext in (".docx", ".xlsx", ".pptx", ".ods", ".odt"):
            bio = BytesIO(file_bytes)
            if not zipfile.is_zipfile(bio):
                # Check for OLE2 header (encrypted office docs)
                if file_bytes[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
                    raise Exception("The file is password-protected and cannot be converted.")
            else:
                try:
                    with zipfile.ZipFile(BytesIO(file_bytes)) as zf:
                        for info in zf.infolist():
                            if info.flag_bits & 0x1:
                                raise Exception("The file contains encrypted content and cannot be converted.")
                except Exception as e:
                    if "password" in str(e).lower() or "encrypted" in str(e).lower():
                        raise

        if ext == ".pdf" and fitz:
            try:
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                is_enc = doc.is_encrypted
                doc.close()
                if is_enc:
                    raise Exception("The PDF file is password-protected and cannot be converted.")
            except Exception as e:
                if "password" in str(e).lower():
                    raise Exception("The PDF file is password-protected and cannot be converted.")

    @staticmethod
    def _get_safe_filename(filename: str) -> str:
        if not filename:
            return None
        return os.path.basename(filename).replace(" ", "_")

    # -----------------------------------------------------------------------
    # Main routing dispatcher
    # -----------------------------------------------------------------------

    @staticmethod
    async def convert(
        file: Union[UploadFile, List[UploadFile]],
        target_format: str = "pdf",
        filename: Optional[str] = None,
    ):
        if filename:
            filename = UniversalConverterController._get_safe_filename(filename)

        target_format = target_format.lower().lstrip(".")

        # Normalise single-item list
        if isinstance(file, list) and len(file) == 1:
            file = file[0]

        if isinstance(file, list):
            return await UniversalConverterController._handle_multiple_images(file, filename)

        ext = os.path.splitext(file.filename)[1].lower()

        # ---- Images ----
        if ext in (".jpg", ".jpeg", ".png", ".bmp", ".gif"):
            return await UniversalConverterController._handle_single_image(file, target_format, filename)

        # ---- HTML ----
        if ext in (".html", ".htm"):
            return await UniversalConverterController._handle_html(file, filename)

        # ---- PDF → Word ----
        if ext == ".pdf" and target_format in ("docx", "doc"):
            base = filename or f"doc_{int(time.time())}_{uuid.uuid4().hex[:8]}.{target_format}"
            if not base.endswith(f".{target_format}"):
                base += f".{target_format}"
            return await UniversalConverterController._handle_pdf_to_word(file, target_format, base)

        # ---- PDF → PDF (passthrough / unsupported) ----
        if ext == ".pdf":
            raise HTTPException(status_code=400, detail="PDF to PDF conversion is not supported.")

        # ---- Unsupported legacy binary formats ----
        if ext in (".doc", ".ppt", ".odp"):
            friendly = {".doc": ".docx", ".ppt": ".pptx", ".odp": ".odp"}
            target = friendly.get(ext, ".docx")
            raise HTTPException(
                status_code=400,
                detail=(
                    f"The '{ext}' format is not supported. "
                    f"Please convert your file to '{target}' and try again."
                ),
            )

        # ---- Office / text formats → PDF ----
        fn = filename or f"doc_{int(time.time())}_{uuid.uuid4().hex[:8]}.pdf"
        if not fn.endswith(".pdf"):
            fn += ".pdf"

        await file.seek(0)
        content = await file.read()
        await file.close()

        UniversalConverterController._validate_file_not_protected(content, ext)

        if ext == ".docx":
            return await UniversalConverterController._handle_docx_to_pdf(content, fn)
        if ext == ".odt":
            return await UniversalConverterController._handle_odt_to_pdf(content, fn)
        if ext == ".rtf":
            return await UniversalConverterController._handle_rtf_to_pdf(content, fn)
        if ext == ".txt":
            return await UniversalConverterController._handle_txt_to_pdf(content, fn)
        if ext == ".xlsx":
            return await UniversalConverterController._handle_xlsx_to_pdf(content, fn)
        if ext == ".xls":
            return await UniversalConverterController._handle_xls_to_pdf(content, fn)
        if ext in (".ods",):
            return await UniversalConverterController._handle_ods_to_pdf(content, fn)
        if ext == ".csv":
            return await UniversalConverterController._handle_csv_to_pdf(content, fn)
        if ext == ".pptx":
            return await UniversalConverterController._handle_pptx_to_pdf(content, fn)

        raise HTTPException(status_code=400, detail=f"Unsupported file format: {ext}")

    # -----------------------------------------------------------------------
    # Image handlers (unchanged)
    # -----------------------------------------------------------------------

    @staticmethod
    async def _handle_single_image(file: UploadFile, target_format: str, filename: str = None):
        if not filename:
            filename = f"img_{int(time.time())}_{uuid.uuid4().hex[:8]}.{target_format}"
        if not filename.endswith(f".{target_format}"):
            filename += f".{target_format}"

        try:
            await file.seek(0)
            content = await file.read()
            if not content:
                raise Exception("Uploaded file is empty")

            img = Image.open(BytesIO(content))
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            elif img.mode != "RGB":
                img = img.convert("RGB")

            output_dir = _make_output_dir()
            output_path = output_dir / filename
            img.save(output_path, "PDF" if target_format.lower() == "pdf" else target_format.upper())

            if not output_path.exists():
                raise Exception("Output file was not saved successfully")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Image conversion failed: {str(e)}")
        finally:
            await file.close()

    @staticmethod
    async def _handle_multiple_images(files: List[UploadFile], filename: str = None):
        if not filename:
            filename = f"images_{int(time.time())}_{uuid.uuid4().hex[:8]}.pdf"

        try:
            images = []
            for f in files:
                await f.seek(0)
                content = await f.read()
                if not content:
                    continue
                img = Image.open(BytesIO(content))
                if img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")
                elif img.mode != "RGB":
                    img = img.convert("RGB")
                images.append(img)

            if not images:
                raise HTTPException(status_code=400, detail="No valid images were provided or identified")

            output_dir = _make_output_dir()
            output_path = output_dir / filename
            images[0].save(output_path, "PDF", resolution=100.0, save_all=True, append_images=images[1:])

            if not output_path.exists():
                raise Exception("Combined PDF was not saved successfully")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Images to PDF failed: {str(e)}")
        finally:
            for f in files:
                await f.close()

    # -----------------------------------------------------------------------
    # HTML handlers (unchanged)
    # -----------------------------------------------------------------------

    @staticmethod
    async def _handle_html(file: UploadFile, filename: str = None):
        if not filename:
            filename = f"html_{int(time.time())}_{uuid.uuid4().hex[:8]}.pdf"
        if not filename.endswith(".pdf"):
            filename += ".pdf"

        output_dir = _make_output_dir()
        output_path = output_dir / filename

        try:
            await file.seek(0)
            content = await file.read()
            html_content = content.decode("utf-8")

            def generate_pdf():
                with open(output_path, "wb") as f:
                    return pisa.CreatePDF(html_content, dest=f)

            pisa_status = await asyncio.to_thread(generate_pdf)
            if pisa_status.err:
                raise Exception(f"xhtml2pdf error: {pisa_status.err}")
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise Exception("PDF file was not created or is empty")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"HTML conversion failed: {str(e)}")
        finally:
            await file.close()

    @staticmethod
    async def convert_html_code(html_code: str, filename: str = None):
        if filename:
            filename = UniversalConverterController._get_safe_filename(filename)
        if not filename:
            filename = f"html_code_{int(time.time())}_{uuid.uuid4().hex[:8]}.pdf"
        if not filename.endswith(".pdf"):
            filename += ".pdf"

        output_dir = _make_output_dir()
        output_path = output_dir / filename

        try:
            def generate_pdf():
                with open(output_path, "wb") as f:
                    return pisa.CreatePDF(html_code, dest=f)

            pisa_status = await asyncio.to_thread(generate_pdf)
            if pisa_status.err:
                raise Exception(f"xhtml2pdf error: {pisa_status.err}")
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise Exception("PDF file was not created or is empty")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"HTML Code conversion failed: {str(e)}")

    @staticmethod
    async def convert_html_link(url: str, filename: str = None):
        if filename:
            filename = UniversalConverterController._get_safe_filename(filename)
        if not filename:
            filename = f"html_link_{int(time.time())}_{uuid.uuid4().hex[:8]}.pdf"
        if not filename.endswith(".pdf"):
            filename += ".pdf"

        output_dir = _make_output_dir()
        output_path = output_dir / filename

        try:
            def fetch_and_generate():
                response = requests.get(url, timeout=15)
                response.raise_for_status()
                with open(output_path, "wb") as f:
                    return pisa.CreatePDF(response.text, dest=f)

            pisa_status = await asyncio.to_thread(fetch_and_generate)
            if pisa_status.err:
                raise Exception(f"xhtml2pdf error: {pisa_status.err}")
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise Exception("PDF file was not created or is empty")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"HTML Link conversion failed: {str(e)}")

    # -----------------------------------------------------------------------
    # PDF → Word (unchanged — pdf2docx is pure Python)
    # -----------------------------------------------------------------------

    @staticmethod
    async def _handle_pdf_to_word(file: UploadFile, target_format: str, filename: str):
        temp_upload_dir = PRIVATE_DISK / "uploads"
        temp_upload_dir.mkdir(parents=True, exist_ok=True)

        ext = os.path.splitext(file.filename)[1].lower()
        temp_input_path = temp_upload_dir / f"{uuid.uuid4().hex}{ext}"
        output_path = _make_output_dir() / filename

        try:
            await file.seek(0)
            content = await file.read()
            with open(temp_input_path, "wb") as buf:
                buf.write(content)

            def perform_conversion():
                cv = Converter(str(temp_input_path))
                cv.convert(str(output_path), start=0, end=None)
                cv.close()

            await asyncio.to_thread(perform_conversion)

            if not output_path.exists():
                raise Exception("Word file was not created by the converter")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PDF to Word conversion failed: {str(e)}")
        finally:
            if os.path.exists(temp_input_path):
                os.remove(temp_input_path)
            await file.close()

    # -----------------------------------------------------------------------
    # DOCX → PDF  (mammoth → HTML → xhtml2pdf)
    # -----------------------------------------------------------------------

    @staticmethod
    async def _handle_docx_to_pdf(content: bytes, filename: str):
        if mammoth is None:
            raise HTTPException(status_code=500, detail="mammoth library is not installed.")

        output_dir = _make_output_dir()
        output_path = output_dir / filename

        try:
            def convert():
                result = mammoth.convert_to_html(BytesIO(content))
                html = _wrap_html(result.value)
                return _html_to_pdf_bytes(html, output_path)

            pisa_status = await asyncio.to_thread(convert)
            if pisa_status.err:
                raise Exception(f"xhtml2pdf error: {pisa_status.err}")
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise Exception("PDF was not created or is empty")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"DOCX to PDF failed: {str(e)}")

    # -----------------------------------------------------------------------
    # TXT → PDF  (reportlab Preformatted)
    # -----------------------------------------------------------------------

    @staticmethod
    async def _handle_txt_to_pdf(content: bytes, filename: str):
        output_dir = _make_output_dir()
        output_path = output_dir / filename

        try:
            def convert():
                text = content.decode("utf-8", errors="replace")
                _rl_text_to_pdf(text, output_path)

            await asyncio.to_thread(convert)
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise Exception("PDF was not created or is empty")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"TXT to PDF failed: {str(e)}")

    # -----------------------------------------------------------------------
    # CSV → PDF  (stdlib csv + xhtml2pdf table)
    # -----------------------------------------------------------------------

    @staticmethod
    async def _handle_csv_to_pdf(content: bytes, filename: str):
        output_dir = _make_output_dir()
        output_path = output_dir / filename

        try:
            def convert():
                text = content.decode("utf-8", errors="replace")
                reader = csv_module.reader(StringIO(text))
                rows = list(reader)
                if not rows:
                    raise Exception("CSV file is empty")
                headers = rows[0]
                data_rows = rows[1:]
                html = _wrap_html(_build_html_table(headers, data_rows, sheet_name="CSV Data"))
                return _html_to_pdf_bytes(html, output_path)

            pisa_status = await asyncio.to_thread(convert)
            if pisa_status.err:
                raise Exception(f"xhtml2pdf error: {pisa_status.err}")
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise Exception("PDF was not created or is empty")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"CSV to PDF failed: {str(e)}")

    # -----------------------------------------------------------------------
    # XLSX → PDF  (openpyxl + xhtml2pdf)
    # -----------------------------------------------------------------------

    @staticmethod
    async def _handle_xlsx_to_pdf(content: bytes, filename: str):
        if openpyxl is None:
            raise HTTPException(status_code=500, detail="openpyxl library is not installed.")

        output_dir = _make_output_dir()
        output_path = output_dir / filename

        try:
            def convert():
                wb = openpyxl.load_workbook(BytesIO(content), data_only=True)
                body = ""
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    all_rows = list(ws.iter_rows(values_only=True))
                    if not all_rows:
                        continue
                    headers = [str(v) if v is not None else "" for v in all_rows[0]]
                    data_rows = [[str(v) if v is not None else "" for v in row] for row in all_rows[1:]]
                    body += _build_html_table(headers, data_rows, sheet_name=sheet_name)
                html = _wrap_html(body)
                return _html_to_pdf_bytes(html, output_path)

            pisa_status = await asyncio.to_thread(convert)
            if pisa_status.err:
                raise Exception(f"xhtml2pdf error: {pisa_status.err}")
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise Exception("PDF was not created or is empty")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"XLSX to PDF failed: {str(e)}")

    # -----------------------------------------------------------------------
    # XLS → PDF  (xlrd + xhtml2pdf)
    # -----------------------------------------------------------------------

    @staticmethod
    async def _handle_xls_to_pdf(content: bytes, filename: str):
        if xlrd is None:
            raise HTTPException(status_code=500, detail="xlrd library is not installed.")

        output_dir = _make_output_dir()
        output_path = output_dir / filename

        try:
            def convert():
                wb = xlrd.open_workbook(file_contents=content)
                body = ""
                for sheet_name in wb.sheet_names():
                    ws = wb.sheet_by_name(sheet_name)
                    if ws.nrows == 0:
                        continue
                    headers = [str(ws.cell_value(0, c)) for c in range(ws.ncols)]
                    data_rows = [
                        [str(ws.cell_value(r, c)) for c in range(ws.ncols)]
                        for r in range(1, ws.nrows)
                    ]
                    body += _build_html_table(headers, data_rows, sheet_name=sheet_name)
                html = _wrap_html(body)
                return _html_to_pdf_bytes(html, output_path)

            pisa_status = await asyncio.to_thread(convert)
            if pisa_status.err:
                raise Exception(f"xhtml2pdf error: {pisa_status.err}")
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise Exception("PDF was not created or is empty")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"XLS to PDF failed: {str(e)}")

    # -----------------------------------------------------------------------
    # ODS → PDF  (odfpy + xhtml2pdf)
    # -----------------------------------------------------------------------

    @staticmethod
    async def _handle_ods_to_pdf(content: bytes, filename: str):
        if odf_load is None:
            raise HTTPException(status_code=500, detail="odfpy library is not installed.")

        output_dir = _make_output_dir()
        output_path = output_dir / filename

        try:
            def convert():
                doc = odf_load(BytesIO(content))
                body = ""
                for sheet in doc.spreadsheet.getElementsByType(odf_table.Table):
                    sheet_name = sheet.getAttribute("name") or "Sheet"
                    rows_data = []
                    for row in sheet.getElementsByType(odf_table.TableRow):
                        cells = []
                        for cell in row.getElementsByType(odf_table.TableCell):
                            ps = cell.getElementsByType(odf_text.P)
                            val = " ".join(
                                "".join(n.data for n in p.childNodes if hasattr(n, "data"))
                                for p in ps
                            )
                            cells.append(val)
                        if any(cells):
                            rows_data.append(cells)
                    if rows_data:
                        headers = rows_data[0]
                        data_rows = rows_data[1:]
                        body += _build_html_table(headers, data_rows, sheet_name=sheet_name)
                html = _wrap_html(body)
                return _html_to_pdf_bytes(html, output_path)

            pisa_status = await asyncio.to_thread(convert)
            if pisa_status.err:
                raise Exception(f"xhtml2pdf error: {pisa_status.err}")
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise Exception("PDF was not created or is empty")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"ODS to PDF failed: {str(e)}")

    # -----------------------------------------------------------------------
    # PPTX → PDF  (python-pptx → slide text → reportlab)
    # -----------------------------------------------------------------------

    @staticmethod
    async def _handle_pptx_to_pdf(content: bytes, filename: str):
        if PptxPresentation is None:
            raise HTTPException(status_code=500, detail="python-pptx library is not installed.")

        output_dir = _make_output_dir()
        output_path = output_dir / filename

        try:
            def convert():
                prs = PptxPresentation(BytesIO(content))
                doc = SimpleDocTemplate(
                    str(output_path),
                    pagesize=A4,
                    rightMargin=20 * mm,
                    leftMargin=20 * mm,
                    topMargin=20 * mm,
                    bottomMargin=20 * mm,
                )
                styles = getSampleStyleSheet()
                story = []

                for i, slide in enumerate(prs.slides, 1):
                    # Slide number heading
                    story.append(Paragraph(f"Slide {i}", styles["Heading1"]))
                    story.append(Spacer(1, 4 * mm))

                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if not text:
                                continue
                            # Use heading style for title shapes
                            if shape.shape_type == 13 or (hasattr(shape, "placeholder_format") and shape.placeholder_format and shape.placeholder_format.idx == 0):
                                story.append(Paragraph(text, styles["Heading2"]))
                            else:
                                story.append(Paragraph(text, styles["Normal"]))

                    story.append(Spacer(1, 8 * mm))

                doc.build(story)

            await asyncio.to_thread(convert)
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise Exception("PDF was not created or is empty")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PPTX to PDF failed: {str(e)}")

    # -----------------------------------------------------------------------
    # RTF → PDF  (striprtf → text → reportlab)
    # -----------------------------------------------------------------------

    @staticmethod
    async def _handle_rtf_to_pdf(content: bytes, filename: str):
        if rtf_to_text is None:
            raise HTTPException(status_code=500, detail="striprtf library is not installed.")

        output_dir = _make_output_dir()
        output_path = output_dir / filename

        try:
            def convert():
                raw = content.decode("utf-8", errors="replace")
                text = rtf_to_text(raw)
                _rl_text_to_pdf(text, output_path)

            await asyncio.to_thread(convert)
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise Exception("PDF was not created or is empty")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"RTF to PDF failed: {str(e)}")

    # -----------------------------------------------------------------------
    # ODT → PDF  (odfpy → plain text → reportlab)
    # -----------------------------------------------------------------------

    @staticmethod
    async def _handle_odt_to_pdf(content: bytes, filename: str):
        if odf_load is None:
            raise HTTPException(status_code=500, detail="odfpy library is not installed.")

        output_dir = _make_output_dir()
        output_path = output_dir / filename

        try:
            def convert():
                doc = odf_load(BytesIO(content))
                text = _extract_odf_text(doc)
                _rl_text_to_pdf(text, output_path)

            await asyncio.to_thread(convert)
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise Exception("PDF was not created or is empty")

            return {"status": "success", "file_url": f"{APP_URL}/storage/converted/{filename}"}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"ODT to PDF failed: {str(e)}")
